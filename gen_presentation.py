"""
Gerador de apresentação PPTX — Tech Challenge Fase 4
FIAP | LSTM Stock Price Prediction
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ─── Paleta ────────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0D, 0x1B, 0x2A)   # azul-escuro quase preto
ACCENT      = RGBColor(0x00, 0xB4, 0xD8)   # ciano vibrante
ACCENT2     = RGBColor(0x48, 0xCA, 0xE4)   # ciano claro
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB0, 0xC4, 0xD8)
DARK_CARD   = RGBColor(0x16, 0x2B, 0x3E)   # card ligeiramente mais claro
GREEN       = RGBColor(0x06, 0xD6, 0x7A)
YELLOW      = RGBColor(0xFF, 0xD1, 0x66)
RED         = RGBColor(0xEF, 0x47, 0x6F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

ARTIFACTS = os.path.join(os.path.dirname(__file__), "artifacts")

# ─── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completamente em branco
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    size=18, color=WHITE, title_color=ACCENT,
                    dot_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  •  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txBox


def add_image(slide, img_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        slide.shapes.add_picture(img_path, left, top, width)


def section_bar(slide, title):
    """Barra lateral esquerda colorida com título vertical."""
    bar = add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, ACCENT)
    add_text(slide, title,
             Inches(0.28), Inches(0.18), Inches(12.7), Inches(0.55),
             size=13, bold=False, color=ACCENT, italic=True)


def slide_title(slide, title, subtitle=None):
    add_text(slide, title,
             Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.75),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.72), Inches(12.6), Inches(0.45),
                 size=16, color=ACCENT)
    # linha divisória
    add_rect(slide, Inches(0.35), Inches(1.18), Inches(12.6), Inches(0.04), ACCENT)


def metric_card(slide, left, top, w, h, label, value, unit="", color=ACCENT):
    add_rect(slide, left, top, w, h, DARK_CARD)
    add_text(slide, label,
             left + Inches(0.15), top + Inches(0.1), w - Inches(0.3), Inches(0.4),
             size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, value,
             left + Inches(0.1), top + Inches(0.42), w - Inches(0.2), Inches(0.65),
             size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    if unit:
        add_text(slide, unit,
                 left + Inches(0.1), top + Inches(1.0), w - Inches(0.2), Inches(0.35),
                 size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_capa(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # gradiente visual — bloco azul-ciano no topo
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)

    add_text(s, "FIAP  |  Tech Challenge  —  Fase 4",
             Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.65),
             size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)

    add_text(s, "Previsão de Preços de Ações",
             Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.4),
             size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "com Redes Neurais LSTM",
             Inches(0.5), Inches(2.9), Inches(12.3), Inches(1.0),
             size=36, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(s, "Apple Inc. (AAPL)  •  2018 – 2026  •  FastAPI  •  ONNX Runtime  •  Vercel",
             Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
             size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text(s, "Deep Learning & IA — 2026",
             Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
             size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    return s


def slide_desafio_solucao(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "01")
    slide_title(s, "O Desafio & A Solução")

    # Coluna esquerda — problema
    add_rect(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.85), DARK_CARD)
    add_text(s, "O PROBLEMA",
             Inches(0.55), Inches(1.5), Inches(5.5), Inches(0.5),
             size=13, bold=True, color=ACCENT)
    items_prob = [
        "Criar pipeline completa de Deep Learning",
        "Modelo LSTM para prever fechamento de ações",
        "Dados históricos → Treino → API → Produção",
        "Monitorar o modelo em produção",
    ]
    add_bullet_list(s, items_prob, Inches(0.55), Inches(2.0),
                    Inches(5.5), Inches(4.5), size=17)

    # Coluna direita — solução
    add_rect(s, Inches(6.55), Inches(1.35), Inches(6.45), Inches(5.85), DARK_CARD)
    add_text(s, "NOSSA SOLUÇÃO",
             Inches(6.75), Inches(1.5), Inches(6.0), Inches(0.5),
             size=13, bold=True, color=GREEN)

    flow = [
        ("📦", "Yahoo Finance (yfinance)", "Coleta de dados AAPL 2018–2026"),
        ("🧠", "LSTM Empilhado", "60 dias × 16 features → delta de preço"),
        ("⚡", "ONNX Runtime", "Conversão para inferência leve em produção"),
        ("🚀", "FastAPI + Vercel", "API RESTful serverless com CI/CD"),
        ("📊", "Streamlit Dashboard", "Monitoramento em tempo real"),
    ]
    y = 2.05
    for icon, title, desc in flow:
        add_text(s, f"{icon}  {title}",
                 Inches(6.75), Inches(y), Inches(6.0), Inches(0.4),
                 size=16, bold=True, color=WHITE)
        add_text(s, f"      {desc}",
                 Inches(6.75), Inches(y + 0.38), Inches(6.0), Inches(0.35),
                 size=13, color=LIGHT_GRAY)
        y += 0.95

    return s


def slide_dados(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "02")
    slide_title(s, "Dados & Engenharia de Features",
                subtitle="Apple Inc. (AAPL) · yfinance · 2018–2026 · 2.011 dias de pregão")

    # Card resumo de dados (esquerda)
    add_rect(s, Inches(0.35), Inches(1.35), Inches(3.5), Inches(5.85), DARK_CARD)
    add_text(s, "DATASET",
             Inches(0.55), Inches(1.5), Inches(3.1), Inches(0.45),
             size=13, bold=True, color=ACCENT)
    info = [
        "Ação: AAPL (Apple Inc.)",
        "Período: 2018-01-01 → 2026-01-01",
        "Total: 2.011 pregões",
        "Treino: 1.393 dias (70%)",
        "Validação: 298 dias (15%)",
        "Teste: 298 dias (15%)",
        "Split: temporal (sem vazamento)",
        "Scaler: RobustScaler",
        "Janela de entrada: 60 dias",
        "Alvo: Δclose (t+1) − close(t)",
    ]
    add_bullet_list(s, info, Inches(0.55), Inches(2.0),
                    Inches(3.1), Inches(5.0), size=14, color=LIGHT_GRAY)

    # Tabela de features (direita)
    add_rect(s, Inches(4.1), Inches(1.35), Inches(8.9), Inches(5.85), DARK_CARD)
    add_text(s, "16 FEATURES UTILIZADAS",
             Inches(4.3), Inches(1.5), Inches(8.4), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    categories = [
        ("OHLCV",       ACCENT,  ["close", "high", "low", "open", "volume"]),
        ("Retornos",    ACCENT2, ["ret_1  (retorno diário %)", "log_ret_1  (log-retorno)"]),
        ("Tendência",   GREEN,   ["sma_7, sma_21  (médias simples)", "ema_12, ema_26  (médias exponenciais)"]),
        ("Momentum",    YELLOW,  ["macd, macd_signal", "rsi_14  (Índice de Força Relativa)"]),
        ("Volatilidade",RED,     ["vol_7, vol_21  (desvio padrão rolling)"]),
    ]

    y = 2.05
    col_w = Inches(4.1)
    for cat, col, feats in categories:
        add_text(s, f"▌ {cat}",
                 Inches(4.3), Inches(y), col_w, Inches(0.38),
                 size=14, bold=True, color=col)
        y += 0.38
        for f in feats:
            add_text(s, f"    {f}",
                     Inches(4.3), Inches(y), Inches(8.5), Inches(0.32),
                     size=13, color=LIGHT_GRAY)
            y += 0.32
        y += 0.08

    return s


def slide_modelo(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "03")
    slide_title(s, "Modelo LSTM",
                subtitle="Long Short-Term Memory · TensorFlow/Keras · Predição de Δ preço")

    # Diagrama da rede (esquerda)
    add_rect(s, Inches(0.35), Inches(1.35), Inches(5.4), Inches(5.85), DARK_CARD)
    add_text(s, "ARQUITETURA",
             Inches(0.55), Inches(1.5), Inches(5.0), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    layers = [
        (ACCENT,  "Input",      "60 timesteps × 16 features"),
        (ACCENT,  "LSTM 1",     "64 unidades · recurrent_dropout=0.05"),
        (ACCENT2, "Dropout",    "0.20"),
        (ACCENT,  "LSTM 2",     "32 unidades · recurrent_dropout=0.05"),
        (ACCENT2, "Dropout",    "0.20"),
        (GREEN,   "Dense",      "16 unidades · ativação ReLU"),
        (YELLOW,  "Output",     "1 unidade → Δclose ($)"),
    ]

    y = 2.1
    for col, name, desc in layers:
        add_rect(s, Inches(0.7), Inches(y), Inches(4.7), Inches(0.52), col)
        add_text(s, f"{name}  —  {desc}",
                 Inches(0.85), Inches(y + 0.06), Inches(4.4), Inches(0.42),
                 size=13, bold=True, color=BG, align=PP_ALIGN.LEFT)
        y += 0.62
        if y < 6.5 and name not in ("Output",):
            add_text(s, "▼",
                     Inches(2.7), Inches(y - 0.14), Inches(0.5), Inches(0.3),
                     size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Hiperparâmetros (direita)
    add_rect(s, Inches(6.05), Inches(1.35), Inches(7.0), Inches(5.85), DARK_CARD)
    add_text(s, "HIPERPARÂMETROS & DECISÕES",
             Inches(6.25), Inches(1.5), Inches(6.6), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    params = [
        ("Otimizador",    "Adam  (lr=5e-4, clipnorm=1.0)"),
        ("Loss",          "Huber Loss  (robusto a outliers)"),
        ("Janela",        "60 dias (lookback)"),
        ("Horizonte",     "1 dia à frente"),
        ("Alvo",          "Δclose = close(t+1) − close(t)"),
        ("Por quê delta?","Série mais estacionária → treino mais estável"),
        ("EarlyStopping", "patience=15  (evita overfitting)"),
        ("ReduceLR",      "patience=8, factor=0.5"),
        ("Split",         "Temporal  —  sem embaralhamento"),
    ]
    y = 2.05
    for label, val in params:
        add_text(s, label,
                 Inches(6.25), Inches(y), Inches(2.5), Inches(0.36),
                 size=13, bold=True, color=LIGHT_GRAY)
        add_text(s, val,
                 Inches(8.75), Inches(y), Inches(4.1), Inches(0.36),
                 size=13, color=WHITE)
        y += 0.52

    return s


def slide_metricas(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "04")
    slide_title(s, "Resultados & Métricas",
                subtitle="Avaliado no conjunto de teste  |  298 dias")

    # Cards de métricas principais
    cards = [
        ("MAE",   "$2,73",  "Erro absoluto médio",    ACCENT),
        ("RMSE",  "$4,09",  "Erro quadrático médio",  ACCENT2),
        ("MAPE",  "1,22%",  "Erro percentual médio",  GREEN),
        ("Dir. Acc.", "47,65%", "Acurácia direcional", YELLOW),
    ]
    cw = Inches(3.0)
    ch = Inches(1.55)
    cx = Inches(0.35)
    for label, val, unit, col in cards:
        metric_card(s, cx, Inches(1.35), cw, ch, label, val, unit, col)
        cx += cw + Inches(0.12)

    # Comparação com baselines
    add_rect(s, Inches(0.35), Inches(3.1), Inches(12.6), Inches(3.4), DARK_CARD)
    add_text(s, "COMPARAÇÃO COM BASELINES",
             Inches(0.55), Inches(3.25), Inches(12.0), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    headers = ["Modelo", "MAE ($)", "RMSE ($)", "MAPE (%)"]
    rows = [
        ("🧠  LSTM (nosso modelo)", "2,73", "4,09", "1,22", GREEN),
        ("📏  Naive (hoje = amanhã)", "2,70", "4,07", "1,21", LIGHT_GRAY),
        ("〰️  SMA-60", "13,54", "16,57", "5,82", RED),
    ]

    # header da tabela
    col_x = [Inches(0.55), Inches(4.5), Inches(7.0), Inches(9.8)]
    col_w = [Inches(3.7), Inches(2.3), Inches(2.3), Inches(2.3)]
    y = 3.72
    for i, h in enumerate(headers):
        add_text(s, h, col_x[i], Inches(y), col_w[i], Inches(0.4),
                 size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    y = 4.15
    for model, mae, rmse, mape, col in rows:
        add_text(s, model, col_x[0], Inches(y), col_w[0], Inches(0.45),
                 size=14, bold=True, color=col)
        for i, val in enumerate([mae, rmse, mape], 1):
            add_text(s, val, col_x[i], Inches(y), col_w[i], Inches(0.45),
                     size=14, color=col, align=PP_ALIGN.CENTER)
        y += 0.5

    add_text(s, "💡  Superar o naive em mercados eficientes é desafiador — o grande ganho está no MAPE 5× menor que o SMA-60.",
             Inches(0.55), Inches(6.3), Inches(12.0), Inches(0.45),
             size=13, color=LIGHT_GRAY, italic=True)

    return s


def slide_graficos(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "05")
    slide_title(s, "Visualizações do Modelo",
                subtitle="Curvas de treino · Preço real vs. previsto · Comparação MAPE")

    img_paths = [
        os.path.join(ARTIFACTS, "training_curves.png"),
        os.path.join(ARTIFACTS, "real_vs_pred_price.png"),
        os.path.join(ARTIFACTS, "mape_comparison.png"),
    ]
    labels = ["Curvas de Treino (Loss)", "Real vs. Previsto (Preço $)", "Comparação MAPE — Modelos"]

    img_w = Inches(4.15)
    x_positions = [Inches(0.35), Inches(4.59), Inches(8.83)]

    for i, (path, label) in enumerate(zip(img_paths, labels)):
        if os.path.exists(path):
            add_image(s, path, x_positions[i], Inches(1.35), img_w, Inches(5.5))
        add_text(s, label,
                 x_positions[i], Inches(6.9), img_w, Inches(0.4),
                 size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    return s


def slide_api(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "06")
    slide_title(s, "API RESTful — FastAPI",
                subtitle="Clean Architecture · Serverless · https://fiap4.vercel.app")

    # Endpoints (esquerda)
    add_rect(s, Inches(0.35), Inches(1.35), Inches(7.2), Inches(5.85), DARK_CARD)
    add_text(s, "ENDPOINTS PRINCIPAIS",
             Inches(0.55), Inches(1.5), Inches(7.0), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    endpoints = [
        ("GET",  "/api/v1/health",                    "Health check & status"),
        ("GET",  "/api/v1/stocks/history",            "Dados históricos AAPL"),
        ("GET",  "/api/v1/stocks/latest?n=30",        "Últimos N pregões"),
        ("POST", "/api/v1/predictions/predict",       "Prever próximos N dias"),
        ("POST", "/api/v1/predictions/predict-custom","Prever com dados próprios"),
        ("GET",  "/api/v1/predictions/model-info",    "Métricas e metadados"),
    ]

    method_color = {
        "GET": GREEN,
        "POST": YELLOW,
    }

    y = 2.05
    for method, path, desc in endpoints:
        add_rect(s, Inches(0.55), Inches(y), Inches(0.75), Inches(0.38),
                 method_color.get(method, ACCENT))
        add_text(s, method,
                 Inches(0.55), Inches(y + 0.03), Inches(0.75), Inches(0.35),
                 size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(s, path,
                 Inches(1.4), Inches(y), Inches(3.8), Inches(0.38),
                 size=13, bold=True, color=WHITE)
        add_text(s, desc,
                 Inches(5.3), Inches(y), Inches(2.0), Inches(0.38),
                 size=12, color=LIGHT_GRAY)
        y += 0.7

    # Exemplo JSON (direita)
    add_rect(s, Inches(7.75), Inches(1.35), Inches(5.25), Inches(5.85), DARK_CARD)
    add_text(s, "EXEMPLO DE RESPOSTA",
             Inches(7.95), Inches(1.5), Inches(4.9), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    json_example = (
        'POST /api/v1/predictions/predict\n'
        '{ "days_ahead": 3 }\n\n'
        '{\n'
        '  "symbol": "AAPL",\n'
        '  "predictions": [\n'
        '    {"date": "2026-02-20",\n'
        '     "predicted_close": 241.25},\n'
        '    {"date": "2026-02-23",\n'
        '     "predicted_close": 242.15},\n'
        '    {"date": "2026-02-24",\n'
        '     "predicted_close": 241.89}\n'
        '  ],\n'
        '  "model_version": "2.0.0",\n'
        '  "metrics": {\n'
        '    "mae": 2.73,\n'
        '    "mape": 1.22\n'
        '  }\n'
        '}'
    )
    txBox = slide.shapes.add_textbox if False else s.shapes.add_textbox(
        Inches(7.95), Inches(2.05), Inches(4.85), Inches(4.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = json_example
    run.font.size = Pt(11)
    run.font.color.rgb = GREEN
    run.font.name = "Courier New"

    return s


def slide_onnx_deploy(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "07")
    slide_title(s, "ONNX Runtime & Deploy na Vercel",
                subtitle="Keras → ONNX · 63% menor · Serverless · CI/CD automático")

    # Comparação Keras vs ONNX
    add_rect(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(2.5), DARK_CARD)
    add_text(s, "KERAS  vs  ONNX RUNTIME",
             Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    compare = [
        ("",           "Keras (.keras)",  "ONNX Runtime"),
        ("Tamanho",    "436 KB",          "160 KB  ✓"),
        ("Dependência","TensorFlow",      "onnxruntime  ✓"),
        ("Cold Start", "Lento",           "Rápido  ✓"),
        ("Portabilidade","Python-only",   "Multi-linguagem  ✓"),
    ]

    col_x = [Inches(0.55), Inches(2.7), Inches(4.55)]
    col_w = [Inches(2.0), Inches(1.7), Inches(2.3)]

    y = 2.05
    for i, row in enumerate(compare):
        header = i == 0
        for j, cell in enumerate(row):
            color = ACCENT if header else (LIGHT_GRAY if j == 0 else (WHITE if j == 1 else GREEN))
            bold = header or j == 2
            add_text(s, cell, col_x[j], Inches(y), col_w[j], Inches(0.45),
                     size=13, bold=bold, color=color)
        y += 0.45

    # Fluxo de conversão
    add_rect(s, Inches(0.35), Inches(4.0), Inches(6.0), Inches(2.85), DARK_CARD)
    add_text(s, "FLUXO DE CONVERSÃO",
             Inches(0.55), Inches(4.15), Inches(5.7), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    steps = [
        ("train_model.py",      "Treina LSTM e salva final_model.keras"),
        ("convert_to_onnx.py",  "Converte via tf2onnx → final_model.onnx"),
        ("API em produção",     "Carrega ONNX — sem TensorFlow instalado"),
    ]
    y = 4.7
    for script, desc in steps:
        add_text(s, f"▶  {script}",
                 Inches(0.55), Inches(y), Inches(2.5), Inches(0.45),
                 size=13, bold=True, color=YELLOW)
        add_text(s, desc,
                 Inches(3.1), Inches(y), Inches(3.0), Inches(0.45),
                 size=13, color=LIGHT_GRAY)
        y += 0.55

    # Deploy Vercel (direita)
    add_rect(s, Inches(6.6), Inches(1.35), Inches(6.42), Inches(5.85), DARK_CARD)
    add_text(s, "DEPLOY — VERCEL",
             Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    deploy_items = [
        ("CI/CD",          "Auto-deploy a cada push na branch main"),
        ("Runtime",        "Python 3.12 · sem GPU"),
        ("Inferência",     "ONNX Runtime (sem TensorFlow)"),
        ("Roteamento",     "vercel.json → api/main.py"),
        ("Escalabilidade", "Serverless — escala automático"),
        ("Ambiente",       "ENV vars via painel Vercel"),
        ("Logs",           "In-memory (sem escrita em disco)"),
        ("URL Pública",    "https://fiap4.vercel.app"),
    ]
    y = 2.1
    for label, val in deploy_items:
        add_text(s, f"{label}:",
                 Inches(6.8), Inches(y), Inches(2.0), Inches(0.42),
                 size=13, bold=True, color=LIGHT_GRAY)
        add_text(s, val,
                 Inches(8.85), Inches(y), Inches(4.0), Inches(0.42),
                 size=13, color=WHITE)
        y += 0.55

    return s


def slide_arquitetura(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "08")
    slide_title(s, "Arquitetura do Projeto",
                subtitle="Clean Architecture · FastAPI · Separação de responsabilidades")

    # Camadas
    layers_info = [
        (ACCENT,   "PRESENTATION",   "Routes · Middlewares · Factories · Dependency Injection"),
        (GREEN,    "DOMAIN",         "Use Cases · Pydantic Models · Interfaces (Repository)"),
        (YELLOW,   "INFRASTRUCTURE", "StockRepositoryImpl (CSV) · Logger · ONNX Runtime"),
    ]

    y = 1.5
    for col, name, desc in layers_info:
        add_rect(s, Inches(0.35), Inches(y), Inches(8.0), Inches(1.2), col)
        add_text(s, name,
                 Inches(0.55), Inches(y + 0.1), Inches(3.5), Inches(0.5),
                 size=18, bold=True, color=BG)
        add_text(s, desc,
                 Inches(4.1), Inches(y + 0.1), Inches(4.0), Inches(1.0),
                 size=14, color=BG)
        y += 1.45

    # Benefícios (direita)
    add_rect(s, Inches(8.6), Inches(1.35), Inches(4.4), Inches(5.85), DARK_CARD)
    add_text(s, "BENEFÍCIOS",
             Inches(8.8), Inches(1.5), Inches(4.0), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    benefits = [
        "Lógica de negócio independente de framework",
        "Fácil substituição do repositório (CSV → DB)",
        "Alta testabilidade por camadas",
        "Middlewares para monitoramento de performance",
        "Injeção de dependência via factories",
        "Logging estruturado (JSON) por ambiente",
    ]
    add_bullet_list(s, benefits, Inches(8.8), Inches(2.05),
                    Inches(4.0), Inches(4.8), size=14, color=LIGHT_GRAY)

    # Fluxo de request
    add_rect(s, Inches(0.35), Inches(5.6), Inches(8.0), Inches(1.6), DARK_CARD)
    add_text(s, "FLUXO DE REQUISIÇÃO",
             Inches(0.55), Inches(5.72), Inches(7.5), Inches(0.35),
             size=12, bold=True, color=ACCENT)
    add_text(s, "HTTP Request → Route → Factory (Dependency Injection) → Use Case → Repository → Response",
             Inches(0.55), Inches(6.1), Inches(7.5), Inches(0.85),
             size=13, color=WHITE)

    return s


def slide_monitoramento(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "09")
    slide_title(s, "Monitoramento — Streamlit Dashboard",
                subtitle="5 páginas · Tempo real · Consome a API em produção")

    pages = [
        (ACCENT,  "1 · Overview",        "Status da API, info do modelo, tendência histórica do AAPL"),
        (GREEN,   "2 · API Metrics",     "Total de requisições, tempo médio de resposta, taxa de erros"),
        (ACCENT2, "3 · Performance",     "Distribuição de latência por endpoint (histograma)"),
        (YELLOW,  "4 · Predictions",     "Interface interativa: selecione N dias → veja gráfico de previsão"),
        (RED,     "5 · Real-time Logs",  "Logs filtráveis por nível (INFO, WARNING, ERROR) com timestamp"),
    ]

    add_rect(s, Inches(0.35), Inches(1.35), Inches(12.6), Inches(4.6), DARK_CARD)

    y = 1.55
    for col, title, desc in pages:
        add_rect(s, Inches(0.55), Inches(y), Inches(0.12), Inches(0.55), col)
        add_text(s, title,
                 Inches(0.78), Inches(y + 0.06), Inches(3.2), Inches(0.45),
                 size=15, bold=True, color=col)
        add_text(s, desc,
                 Inches(4.1), Inches(y + 0.06), Inches(8.5), Inches(0.45),
                 size=14, color=LIGHT_GRAY)
        y += 0.78

    add_rect(s, Inches(0.35), Inches(6.05), Inches(12.6), Inches(1.2), DARK_CARD)
    add_text(s, "PERFORMANCE MIDDLEWARE",
             Inches(0.55), Inches(6.15), Inches(5.0), Inches(0.35),
             size=12, bold=True, color=ACCENT)
    details = [
        "• Captura métricas de cada request/response automaticamente",
        "• Headers customizados: X-Response-Time · X-Request-ID",
        "• Logs em JSON (local: arquivo · produção: in-memory)",
    ]
    add_bullet_list(s, details, Inches(0.55), Inches(6.52),
                    Inches(12.0), Inches(0.65), size=13, color=LIGHT_GRAY)

    return s


def slide_conclusao(prs):
    s = blank_slide(prs)
    fill_bg(s)
    section_bar(s, "10")
    slide_title(s, "Conclusões & Aprendizados")

    # Coluna esquerda — entregáveis
    add_rect(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(5.85), DARK_CARD)
    add_text(s, "ENTREGÁVEIS ✅",
             Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.45),
             size=13, bold=True, color=GREEN)

    entregaveis = [
        "Modelo LSTM treinado e salvo (Keras + ONNX)",
        "Pipeline completa: coleta → treino → deploy",
        "API RESTful (FastAPI) em produção na Vercel",
        "Dashboard de monitoramento (Streamlit)",
        "Código-fonte no GitHub (Clean Architecture)",
        "16 features com indicadores técnicos",
        "Métricas: MAE $2,73 · MAPE 1,22%",
    ]
    add_bullet_list(s, entregaveis, Inches(0.55), Inches(2.05),
                    Inches(5.6), Inches(4.8), size=16, color=WHITE)

    # Coluna direita — aprendizados
    add_rect(s, Inches(6.6), Inches(1.35), Inches(6.42), Inches(5.85), DARK_CARD)
    add_text(s, "APRENDIZADOS",
             Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.45),
             size=13, bold=True, color=ACCENT)

    aprendizados = [
        "Mercados eficientes → difícil superar baseline naive",
        "Engenharia de features é tão importante quanto a rede",
        "ONNX é essencial para deploys serverless leves",
        "Split temporal evita vazamento de dados futuros",
        "Clean Architecture facilita manutenção e testes",
        "Monitoramento é parte do produto, não opcional",
    ]
    add_bullet_list(s, aprendizados, Inches(6.8), Inches(2.05),
                    Inches(6.0), Inches(4.8), size=16, color=WHITE)

    return s


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_capa(prs)
    slide_desafio_solucao(prs)
    slide_dados(prs)
    slide_modelo(prs)
    slide_metricas(prs)
    slide_graficos(prs)
    slide_api(prs)
    slide_onnx_deploy(prs)
    slide_arquitetura(prs)
    slide_monitoramento(prs)
    slide_conclusao(prs)

    out = os.path.join(os.path.dirname(__file__), "apresentacao_tech_challenge.pptx")
    prs.save(out)
    print(f"✅  Apresentação salva em: {out}")
    print(f"    Slides: {len(prs.slides)}  |  Tempo estimado: ~10 min")


if __name__ == "__main__":
    main()
